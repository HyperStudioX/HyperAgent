"""PPTX generation service for creating professional presentations."""

import base64
import urllib.request
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

from app.core.logging import get_logger

logger = get_logger(__name__)


# ---------------------------------------------------------------------------
# Pydantic models
# ---------------------------------------------------------------------------


class SlideTheme(BaseModel):
    """Visual theme for the slide deck."""

    primary_color: str = "#1B365D"
    secondary_color: str = "#4A90D9"
    accent_color: str = "#E8792F"
    font_heading: str = "Calibri"
    font_body: str = "Calibri"
    background_color: str = "#FFFFFF"


class SlideElement(BaseModel):
    """A single element on a slide."""

    type: str  # "text", "image", "shape"
    content: str = ""  # text content, image URL, or base64 data
    position: dict[str, float] = Field(
        default_factory=lambda: {"left": 0.5, "top": 1.5, "width": 9.0, "height": 5.0}
    )
    style: dict[str, Any] = Field(default_factory=dict)


class SlideSpec(BaseModel):
    """Specification for a single slide."""

    layout: str = "content"  # title_slide, section_header, content, two_column, image_focus, blank
    title: str = ""
    subtitle: str | None = None
    elements: list[SlideElement] = Field(default_factory=list)
    notes: str | None = None
    background_color: str | None = None


class SlideDeck(BaseModel):
    """Complete slide deck specification."""

    title: str
    author: str | None = None
    theme: SlideTheme = Field(default_factory=SlideTheme)
    slides: list[SlideSpec] = Field(default_factory=list)


# ---------------------------------------------------------------------------
# Alignment mapping
# ---------------------------------------------------------------------------

_ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _parse_hex_color(hex_str: str) -> RGBColor:
    """Parse a hex color string like '#FF5733' into an RGBColor."""
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        hex_str = "000000"
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Service
# ---------------------------------------------------------------------------


class PptxGenerationService:
    """Service for generating PPTX presentations from structured specifications."""

    # Widescreen 16:9 dimensions
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def generate_pptx(self, deck: SlideDeck) -> bytes:
        """Generate a PPTX file from a SlideDeck specification.

        Args:
            deck: The slide deck specification.

        Returns:
            Raw PPTX file bytes.
        """
        prs = Presentation()
        self._apply_theme(prs, deck.theme)

        for idx, slide_spec in enumerate(deck.slides):
            layout = slide_spec.layout.lower().replace(" ", "_")
            try:
                if layout == "title_slide":
                    slide = self._create_title_slide(prs, slide_spec, deck.theme)
                elif layout == "section_header":
                    slide = self._create_section_header_slide(prs, slide_spec, deck.theme)
                elif layout == "two_column":
                    slide = self._create_two_column_slide(prs, slide_spec, deck.theme)
                elif layout == "image_focus":
                    slide = self._create_image_slide(prs, slide_spec, deck.theme)
                elif layout == "blank":
                    slide = self._create_blank_slide(prs, slide_spec, deck.theme)
                else:
                    # Default: content layout
                    slide = self._create_content_slide(prs, slide_spec, deck.theme)

                # Add speaker notes
                if slide_spec.notes and slide.has_notes_slide:
                    slide.notes_slide.notes_text_frame.text = slide_spec.notes

                # Add slide number (skip title slide)
                if layout != "title_slide":
                    self._add_slide_number(slide, prs, deck.theme, idx + 1)

            except Exception as e:
                logger.warning("slide_creation_failed", slide_index=idx, error=str(e))
                # Create a fallback blank slide so the deck still has the right count
                blank_layout = prs.slide_layouts[6]  # blank
                prs.slides.add_slide(blank_layout)

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    # ------------------------------------------------------------------
    # Theme
    # ------------------------------------------------------------------

    def _apply_theme(self, prs: Presentation, theme: SlideTheme) -> None:
        """Set slide dimensions to 16:9 widescreen."""
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

    # ------------------------------------------------------------------
    # Slide builders
    # ------------------------------------------------------------------

    def _create_title_slide(self, prs: Presentation, spec: SlideSpec, theme: SlideTheme):
        """Create a title slide with centered title and subtitle."""
        slide_layout = prs.slide_layouts[6]  # blank for full control
        slide = prs.slides.add_slide(slide_layout)

        # Background
        self._set_slide_background(slide, spec.background_color or theme.background_color)

        # Title - large centered text
        left = Inches(1.0)
        top = Inches(2.2)
        width = Inches(11.333)
        height = Inches(1.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = spec.title
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = _parse_hex_color(theme.primary_color)
        run.font.name = theme.font_heading

        # Subtitle
        if spec.subtitle:
            sub_top = Inches(4.2)
            sub_height = Inches(1.2)
            txBox2 = slide.shapes.add_textbox(left, sub_top, width, sub_height)
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = spec.subtitle
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.runs[0]
            run2.font.size = Pt(22)
            run2.font.color.rgb = _parse_hex_color(theme.secondary_color)
            run2.font.name = theme.font_body

        # Accent line below title
        line_left = Inches(4.5)
        line_top = Inches(4.0)
        line_width = Inches(4.333)
        line_height = Inches(0.04)
        shape = slide.shapes.add_shape(1, line_left, line_top, line_width, line_height)  # rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_hex_color(theme.accent_color)
        shape.line.fill.background()

        return slide

    def _create_section_header_slide(self, prs: Presentation, spec: SlideSpec, theme: SlideTheme):
        """Create a section header / divider slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, spec.background_color or theme.primary_color)

        # Section title in white
        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(10.333)
        height = Inches(2.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = spec.title
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = theme.font_heading

        if spec.subtitle:
            p2 = tf.add_paragraph()
            p2.text = spec.subtitle
            p2.alignment = PP_ALIGN.LEFT
            run2 = p2.add_run()
            run2.text = spec.subtitle
            # Fix: clear the auto-created text, use the run instead
            p2.clear()
            r2 = p2.add_run()
            r2.text = spec.subtitle
            r2.font.size = Pt(20)
            r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            r2.font.name = theme.font_body

        return slide

    def _create_content_slide(self, prs: Presentation, spec: SlideSpec, theme: SlideTheme):
        """Create a standard content slide with title and bullet points.

        If the slide contains both text and image elements, delegates to
        _create_content_with_image_slide for a split layout.
        """
        # Check if we have a mix of text and image elements
        has_text = any(el.type == "text" for el in spec.elements)
        has_image = any(el.type == "image" for el in spec.elements)
        if has_text and has_image:
            return self._create_content_with_image_slide(prs, spec, theme)

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, spec.background_color or theme.background_color)

        # Title bar
        title_left = Inches(0.7)
        title_top = Inches(0.4)
        title_width = Inches(11.933)
        title_height = Inches(0.9)
        txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = spec.title
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = _parse_hex_color(theme.primary_color)
        run.font.name = theme.font_heading

        # Accent underline
        line_shape = slide.shapes.add_shape(1, title_left, Inches(1.3), Inches(1.5), Inches(0.03))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = _parse_hex_color(theme.accent_color)
        line_shape.line.fill.background()

        # Content area - bullet points from elements
        content_left = Inches(0.7)
        content_top = Inches(1.6)
        content_width = Inches(11.933)
        content_height = Inches(5.2)
        txBox2 = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        first = True
        for element in spec.elements:
            if element.type == "text":
                if first:
                    p = tf2.paragraphs[0]
                    first = False
                else:
                    p = tf2.add_paragraph()
                p.text = f"\u2022 {element.content}"
                p.space_after = Pt(12)
                self._apply_text_style(p, element.style, theme, default_font_size=18)
            elif element.type == "image":
                self._embed_image(slide, element)

        return slide

    def _create_content_with_image_slide(
        self, prs: Presentation, spec: SlideSpec, theme: SlideTheme
    ):
        """Create a content slide with text on the left (60%) and image on the right (40%)."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, spec.background_color or theme.background_color)

        # Title bar (full width)
        title_left = Inches(0.7)
        title_top = Inches(0.4)
        title_width = Inches(11.933)
        title_height = Inches(0.9)
        txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = spec.title
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = _parse_hex_color(theme.primary_color)
        run.font.name = theme.font_heading

        # Accent underline
        line_shape = slide.shapes.add_shape(1, title_left, Inches(1.3), Inches(1.5), Inches(0.03))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = _parse_hex_color(theme.accent_color)
        line_shape.line.fill.background()

        # Left column: text (60% width)
        text_left = Inches(0.7)
        text_top = Inches(1.6)
        text_width = Inches(7.0)
        text_height = Inches(5.2)
        txBox2 = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        first = True
        for element in spec.elements:
            if element.type == "text":
                if first:
                    p = tf2.paragraphs[0]
                    first = False
                else:
                    p = tf2.add_paragraph()
                p.text = f"\u2022 {element.content}"
                p.space_after = Pt(12)
                self._apply_text_style(p, element.style, theme, default_font_size=18)

        # Right column: first image (40% width)
        for element in spec.elements:
            if element.type == "image":
                element.position = {
                    "left": 8.0,
                    "top": 1.6,
                    "width": 4.633,
                    "height": 5.2,
                }
                self._embed_image(slide, element)
                break

        return slide

    def _create_two_column_slide(self, prs: Presentation, spec: SlideSpec, theme: SlideTheme):
        """Create a two-column content slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, spec.background_color or theme.background_color)

        # Title
        title_left = Inches(0.7)
        title_top = Inches(0.4)
        title_width = Inches(11.933)
        title_height = Inches(0.9)
        txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = spec.title
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = _parse_hex_color(theme.primary_color)
        run.font.name = theme.font_heading

        # Split elements roughly in half for two columns
        elements = spec.elements
        mid = len(elements) // 2 if len(elements) > 1 else len(elements)
        left_elements = elements[:mid]
        right_elements = elements[mid:]

        # Left column
        col_top = Inches(1.6)
        col_height = Inches(5.2)
        left_box = slide.shapes.add_textbox(Inches(0.7), col_top, Inches(5.8), col_height)
        ltf = left_box.text_frame
        ltf.word_wrap = True
        first = True
        for el in left_elements:
            if el.type == "text":
                if first:
                    lp = ltf.paragraphs[0]
                    first = False
                else:
                    lp = ltf.add_paragraph()
                lp.text = el.content
                lp.space_after = Pt(6)
                self._apply_text_style(lp, el.style, theme)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.833), col_top, Inches(5.8), col_height)
        rtf = right_box.text_frame
        rtf.word_wrap = True
        first = True
        for el in right_elements:
            if el.type == "text":
                if first:
                    rp = rtf.paragraphs[0]
                    first = False
                else:
                    rp = rtf.add_paragraph()
                rp.text = el.content
                rp.space_after = Pt(6)
                self._apply_text_style(rp, el.style, theme)

        return slide

    def _create_image_slide(self, prs: Presentation, spec: SlideSpec, theme: SlideTheme):
        """Create an image-focused slide with optional caption."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, spec.background_color or theme.background_color)

        # Title (small)
        if spec.title:
            txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(11.933), Inches(0.7))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = spec.title
            run = p.runs[0]
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = _parse_hex_color(theme.primary_color)
            run.font.name = theme.font_heading

        # Try to find and embed the first image element
        for element in spec.elements:
            if element.type == "image":
                # Large centered image
                element.position = {"left": 1.5, "top": 1.2, "width": 10.333, "height": 5.5}
                self._embed_image(slide, element)
                break
            elif element.type == "text":
                # Treat as caption below image area
                txBox2 = slide.shapes.add_textbox(
                    Inches(1.5), Inches(6.8), Inches(10.333), Inches(0.5)
                )
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = element.content
                p2.alignment = PP_ALIGN.CENTER
                self._apply_text_style(p2, element.style, theme)

        return slide

    def _create_blank_slide(self, prs: Presentation, spec: SlideSpec, theme: SlideTheme):
        """Create a blank slide with custom-positioned elements."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, spec.background_color or theme.background_color)

        for element in spec.elements:
            pos = element.position
            left = Inches(pos.get("left", 0.5))
            top = Inches(pos.get("top", 0.5))
            width = Inches(pos.get("width", 5.0))
            height = Inches(pos.get("height", 3.0))

            if element.type == "text":
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = element.content
                self._apply_text_style(p, element.style, theme)
            elif element.type == "image":
                self._embed_image(slide, element)

        return slide

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _add_slide_number(self, slide, prs: Presentation, theme: SlideTheme, number: int) -> None:
        """Add a slide number in the bottom-right corner."""
        left = Inches(12.0)
        top = Inches(7.0)
        width = Inches(1.0)
        height = Inches(0.35)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.alignment = PP_ALIGN.RIGHT
        run = p.runs[0]
        run.font.size = Pt(10)
        run.font.color.rgb = _parse_hex_color(theme.secondary_color)
        run.font.name = theme.font_body

    def _set_slide_background(self, slide, color_hex: str) -> None:
        """Set a solid background color for a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _parse_hex_color(color_hex)

    def _embed_image(self, slide, element: SlideElement) -> None:
        """Download/decode an image and embed it on the slide."""
        pos = element.position
        left = Inches(pos.get("left", 1.0))
        top = Inches(pos.get("top", 1.5))
        width = Inches(pos.get("width", 8.0))
        height = Inches(pos.get("height", 4.5))

        content = element.content.strip()
        image_stream: BytesIO | None = None

        try:
            if content.startswith("data:"):
                # data URI: data:image/png;base64,...
                _, encoded = content.split(",", 1)
                image_stream = BytesIO(base64.b64decode(encoded))
            elif content.startswith("http://") or content.startswith("https://"):
                # URL - download synchronously
                req = urllib.request.Request(content, headers={"User-Agent": "HyperAgent/1.0"})
                with urllib.request.urlopen(req, timeout=15) as resp:
                    image_stream = BytesIO(resp.read())
            else:
                # Assume raw base64
                try:
                    image_stream = BytesIO(base64.b64decode(content))
                except Exception:
                    logger.warning("image_decode_failed", content_preview=content[:60])
                    return

            if image_stream:
                slide.shapes.add_picture(image_stream, left, top, width, height)
        except Exception as e:
            logger.warning("image_embed_failed", error=str(e))

    def _apply_text_style(
        self,
        paragraph,
        style_dict: dict[str, Any],
        theme: SlideTheme,
        default_font_size: int = 18,
    ) -> None:
        """Apply styling to a paragraph's runs."""
        font_size = style_dict.get("font_size")
        bold = style_dict.get("bold")
        italic = style_dict.get("italic")
        color = style_dict.get("color")
        alignment = style_dict.get("alignment")

        if alignment and alignment in _ALIGNMENT_MAP:
            paragraph.alignment = _ALIGNMENT_MAP[alignment]

        fallback_size = Pt(int(font_size)) if font_size else Pt(default_font_size)

        # Apply to existing runs
        for run in paragraph.runs:
            run.font.name = style_dict.get("font", theme.font_body)
            run.font.size = fallback_size
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
            if color:
                run.font.color.rgb = _parse_hex_color(color)
            else:
                run.font.color.rgb = _parse_hex_color(theme.primary_color)

        # If no runs exist (empty paragraph that was just created), set defaults
        if not paragraph.runs and paragraph.text:
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            if not run.text and paragraph.text:
                # Text was set on paragraph directly, re-apply on run
                pass
            run.font.name = style_dict.get("font", theme.font_body)
            run.font.size = fallback_size
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
            if color:
                run.font.color.rgb = _parse_hex_color(color)
            else:
                run.font.color.rgb = _parse_hex_color(theme.primary_color)


# Module-level instance
pptx_generation_service = PptxGenerationService()
